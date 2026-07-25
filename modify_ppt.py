"""
modify_ppt.py — Modifies IDEA_Presentation_Format.pptx according to hackathon instructions:
1. Deletes Slide 1 (Instructions slide)
2. Fills Slide 1 (Title Slide) with student details and problem statement
3. Fills Slide 2 (Proposed Solution) with CyberShield AI architecture & uniqueness
4. Fills Slide 3 (Technical Approach) with tech stack & end-to-end methodology
5. Fills Slide 4 (Feasibility & Viability) with quantitative metrics & risk mitigation
6. Fills Slide 5 (Artifacts) with GitHub repo link, dashboard overview & components
7. Fills Slide 6 (Research & References) with academic references & links
"""

import sys
from pathlib import Path
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = Path("IDEA_Presentation_Format.pptx")
OUTPUT_PPT_PATH = Path("IDEA_Presentation_Format.pptx")


def delete_slide(prs, index):
    """Deletes a slide by index from presentation."""
    rId = prs.slides._sldIdLst[index].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def update_textbox_text(shape, new_text, font_size=14, bold=False):
    """Safely updates text inside a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    # Clear existing paragraphs except first
    p = tf.paragraphs[0]
    p.text = new_text
    p.font.size = Pt(font_size)
    p.font.name = "Helvetica"
    p.font.bold = bold


def populate_presentation():
    prs = pptx.Presentation(PPT_PATH)
    print(f"Initial slides count: {len(prs.slides)}")

    # Delete slide 1 (Instructions slide) if present
    if len(prs.slides) == 7:
        delete_slide(prs, 0)
        print("Deleted Instruction Slide (Slide 1). New count:", len(prs.slides))

    # ════════════════════════════════════════════════════════════
    # SLIDE 1 (TITLE PAGE)
    # ════════════════════════════════════════════════════════════
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if "Problem Statement ID" in text or "Student Name" in text:
                shape.text_frame.text = (
                    "Problem Statement ID: HW-BAD-2026\n"
                    "Problem Statement Title: AI-Powered Behavioral Anomaly Detection for Cybersecurity\n"
                    "Theme: AI/ML & Cybersecurity System Architecture\n"
                    "PS Category: Software\n"
                    "Student Name: Ansh Singh\n"
                    "Student ID: ANSH-SIH-2026\n"
                    "GitHub Repository: https://github.com/SinghAnsh07/cybershield-AI"
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(15)
                    p.font.name = "Helvetica"
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(15, 23, 42)

    # ════════════════════════════════════════════════════════════
    # SLIDE 2 (PROPOSED SOLUTION / IDEA TITLE)
    # ════════════════════════════════════════════════════════════
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            if "IDEA TITLE" in shape.text_frame.text:
                shape.text_frame.text = "IDEA TITLE: CyberShield AI"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.bold = True
            elif "Proposed Solution" in shape.text_frame.text:
                shape.text_frame.text = (
                    "PROPOSED SOLUTION: Behavioral Anomaly Detection & Attack Classification\n\n"
                    "* Multi-Stage Anomaly Engine: Combines a Statistical Profiler (Mahalanobis Distance), "
                    "a PyTorch Autoencoder, a 2-layer BiLSTM sequence detector with Focal Loss, and an XGBoost multi-class attack classifier.\n"
                    "* Solves Core Security Challenges: Handles class imbalance (<0.5% attacks via Focal Loss & SMOTE), "
                    "concept drift (exponential age-decay sample weighting), and cold-start entities (cohort baseline blending).\n"
                    "* Plain-English Explainability: Integrates SHAP TreeExplainer to deliver human-readable alert reasons "
                    "(e.g., 'Flagged due to impossible geo-velocity + novel device fingerprint') for fast SOC triage.\n"
                    "* Interactive SOC Dashboard: Streamlit analyst-facing application with KPI overview, ranked alert queues, "
                    "and entity investigation timelines."
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(13)

    # ════════════════════════════════════════════════════════════
    # SLIDE 3 (TECHNICAL APPROACH)
    # ════════════════════════════════════════════════════════════
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            if "TECHNICAL APPROACH" in shape.text_frame.text:
                shape.text_frame.text = "TECHNICAL APPROACH & METHODOLOGY"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.bold = True
            elif "Technologies to be used" in shape.text_frame.text:
                shape.text_frame.text = (
                    "TECHNOLOGY STACK & PIPELINE FLOW\n\n"
                    "* Programming & Core Libraries: Python 3.12, NumPy, Pandas, Faker (Synthetic Generator), Scikit-Learn.\n"
                    "* Deep Learning & ML: PyTorch (Autoencoder & BiLSTM Detector), XGBoost (Multi-Class Classifier), Imbalanced-Learn (SMOTE).\n"
                    "* Explainability & UI: SHAP TreeExplainer, Streamlit Dashboard, Plotly Visualization Engine, FPDF2.\n"
                    "* End-to-End Execution Flow:\n"
                    "  1. Vectorized Log Ingestion -> 10-Field Telemetry Schema across Users, Service Accounts & Edge Devices.\n"
                    "  2. Feature Extraction Pipeline -> 26 behavioral features (geo-velocity, resource novelty, failed auth counts).\n"
                    "  3. Hybrid Profiling -> Autoencoder reconstruction MSE + Mahalanobis distance scoring.\n"
                    "  4. Sequence Anomaly Scoring -> 10-event window BiLSTM with Focal Loss.\n"
                    "  5. Attack Classification & Explanation -> XGBoost + SHAP natural language reason translation.\n"
                    "* Streaming Blueprint: Apache Kafka + Apache Flink stateful windowing in RocksDB & ONNX sub-50ms serving."
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(12)

    # ════════════════════════════════════════════════════════════
    # SLIDE 4 (FEASIBILITY AND VIABILITY)
    # ════════════════════════════════════════════════════════════
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            if "FEASIBILITY AND VIABILITY" in shape.text_frame.text:
                shape.text_frame.text = "FEASIBILITY, RESULTS & RISK MITIGATION"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.bold = True
            elif "Analysis of the feasibility" in shape.text_frame.text:
                shape.text_frame.text = (
                    "EXPERIMENTAL RESULTS & OPERATIONAL VIABILITY\n\n"
                    "* Empirical Validation: Successfully evaluated on 212,475 access log events across 500 entities over 90 days.\n"
                    "* Low False Positive Rate: Achieved FPR = 0.91% at top 1% analyst alert budget (maintains strict SOC operational compliance).\n"
                    "* High Multi-Class Accuracy: Multi-Class Weighted F1 Score = 0.9117 across 7 injected attack types.\n"
                    "* Key Risk Mitigations:\n"
                    "  - Class Imbalance (<0.5% rare attacks) -> Focal Loss (gamma=2.0) + SMOTE minority oversampling.\n"
                    "  - Cold-Start Entities -> Cohort baseline blending smoothly transitions to individual profiles over 50 events.\n"
                    "  - Concept Drift -> Exponential age-decay weighting (0.95^t) prioritizes recent behavioral trends.\n"
                    "  - High Event Throughput -> Kafka/Flink sub-50ms design with embedded ONNX runtime model scoring."
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(12)

    # ════════════════════════════════════════════════════════════
    # SLIDE 5 (ARTIFACTS)
    # ════════════════════════════════════════════════════════════
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            if "ARTIFACTS" in shape.text_frame.text:
                shape.text_frame.text = "SOLUTION ARTIFACTS & PROTOTYPE"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.bold = True
            elif "Relevant artifacts" in shape.text_frame.text:
                shape.text_frame.text = (
                    "COMPLETE PRODUCTION DELIVERABLES\n\n"
                    "* GitHub Codebase: https://github.com/SinghAnsh07/cybershield-AI (Fully modular open-source implementation)\n"
                    "* End-to-End Pipeline: Single-command execution (`python pipeline.py`) for data generation, feature engineering, "
                    "Autoencoder/BiLSTM/XGBoost training, evaluation, and report compilation.\n"
                    "* SOC Analyst Dashboard: Streamlit interactive UI (`streamlit run dashboard/app.py`) with risk metrics, alert queue, "
                    "SHAP waterfall plots, and entity timelines.\n"
                    "* Production PDF Report: 5-page publication-quality technical document (`report.pdf`) with embedded ROC/PR curves.\n"
                    "* Trained Model Artifacts: Saved weights in `data/models/` (`baseline_profiler.pt`, `detection_model.pt`, `anomaly_classifier.joblib`)."
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(12)

    # ════════════════════════════════════════════════════════════
    # SLIDE 6 (RESEARCH AND REFERENCES)
    # ════════════════════════════════════════════════════════════
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            if "RESEARCH" in shape.text_frame.text:
                shape.text_frame.text = "RESEARCH & REFERENCES"
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.bold = True
            elif "Details / Links" in shape.text_frame.text:
                shape.text_frame.text = (
                    "KEY REFERENCES & RESEARCH WORK\n\n"
                    "1. User and Entity Behavior Analytics (UEBA): IEEE Transactions on Information Forensics and Security on access log anomaly detection.\n"
                    "2. Sequence Anomaly Detection: PyTorch Bidirectional LSTM with Focal Loss for imbalanced temporal security events.\n"
                    "3. Model Explainability: Lundberg & Lee (2017), 'A Unified Approach to Interpreting Model Predictions' (SHAP TreeExplainer).\n"
                    "4. Real-Time Streaming Systems: Apache Flink Stateful Stream Processing & RocksDB state backend architecture.\n"
                    "5. Solution Source Code & Deliverables: https://github.com/SinghAnsh07/cybershield-AI"
                )
                for p in shape.text_frame.paragraphs:
                    p.font.size = Pt(13)

    prs.save(OUTPUT_PPT_PATH)
    print(f"Successfully populated and saved PPT to: {OUTPUT_PPT_PATH}")


if __name__ == "__main__":
    populate_presentation()
